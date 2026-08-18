"""
Transformer Seminar HTML → PDF & PowerPoint(PPTX) 자동 생성 스크립트
1. Playwright(MS Edge)를 통해 HTML 47개 슬라이드 전체 및 수식 요소(KaTeX)를 고해상도로 캡처
2. 고품질 PDF (transformer_seminar.pdf) 생성
3. 수식 캡처 이미지 삽입형 PPTX (transformer_seminar.pptx) 생성
4. 프레젠테이션용 전체 뷰 PPTX (transformer_seminar_presentation.pptx) 생성
"""

import io
import os
import re
import sys
from pathlib import Path
from bs4 import BeautifulSoup
from PIL import Image
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# Windows 콘솔 UTF-8 출력 강제 설정
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")

BASE_DIR = Path(r"c:\python\AI_Seminar")
HTML_FILE = BASE_DIR / "transformer_seminar_new.html"
CAPTURES_DIR = BASE_DIR / "captures"
SLIDES_IMG_DIR = CAPTURES_DIR / "slides"
MATH_IMG_DIR = CAPTURES_DIR / "math"

PDF_OUTPUT = BASE_DIR / "transformer_seminar.pdf"
PPTX_HYBRID_OUTPUT = BASE_DIR / "transformer_seminar_hybrid.pptx"
PPTX_FULL_OUTPUT = BASE_DIR / "transformer_seminar_presentation.pptx"
PPTX_DEFAULT_OUTPUT = BASE_DIR / "transformer_seminar.pptx"

def safe_save_prs(prs, primary_path: Path, fallback_path: Path):
    try:
        prs.save(str(primary_path))
        print(f"[+] PPTX 저장 완료: {primary_path}")
        return primary_path
    except PermissionError:
        print(f"[!] {primary_path.name} 파일이 다른 프로그램(예: 파워포인트)에서 열려 있어 {fallback_path.name} 로 저장합니다.")
        prs.save(str(fallback_path))
        print(f"[+] PPTX 저장 완료: {fallback_path}")
        return fallback_path


# ─── 색상 상수 ─────────────────────────────────────────────────────────────
BG_DARK       = RGBColor(0x07, 0x0a, 0x12)   # 배경 (#070a12)
CARD_BG       = RGBColor(0x12, 0x1a, 0x2b)   # 카드 배경
Q_COLOR       = RGBColor(0x38, 0xbd, 0xf8)   # Query - Cyan
K_COLOR       = RGBColor(0x34, 0xd3, 0x99)   # Key - Emerald
V_COLOR       = RGBColor(0xc0, 0x84, 0xfc)   # Value - Purple
SCORE_COLOR   = RGBColor(0xfb, 0xbf, 0x24)   # Score/Attention - Amber
HW_COLOR      = RGBColor(0xf4, 0x3f, 0x5e)   # H/W Memory - Rose
TEXT_MAIN     = RGBColor(0xf8, 0xfa, 0xfc)   # 메인 텍스트
TEXT_SUB      = RGBColor(0x94, 0xa3, 0xb8)   # 서브 텍스트
HUMAN_NOTE_FG = RGBColor(0xfd, 0xe6, 0x8a)   # 문과 직관 박스 텍스트
BORDER_DARK   = RGBColor(0x25, 0x35, 0x4d)   # 테두리 다크

CARD_COLORS = {
    "card-q":     (Q_COLOR,     RGBColor(0x04, 0x13, 0x1e)),
    "card-k":     (K_COLOR,     RGBColor(0x04, 0x1c, 0x14)),
    "card-v":     (V_COLOR,     RGBColor(0x14, 0x09, 0x1e)),
    "card-score": (SCORE_COLOR, RGBColor(0x1e, 0x13, 0x04)),
    "card-hw":    (HW_COLOR,    RGBColor(0x1e, 0x04, 0x08)),
}

TAG_COLORS = {
    "tag-concept": Q_COLOR,
    "tag-math":    K_COLOR,
    "tag-hw":      HW_COLOR,
}

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def clean_text(el) -> str:
    if el is None:
        return ""
    text = el.get_text(separator=" ", strip=True)
    return re.sub(r"\s+", " ", text).strip()


def add_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK


def add_rect_shape(slide, left, top, width, height, bg_color, border_color=None, border_pt=1.0):
    sp = slide.shapes.add_shape(1, left, top, width, height)
    sp.fill.solid()
    sp.fill.fore_color.rgb = bg_color
    if border_color:
        sp.line.color.rgb = border_color
        sp.line.width = Pt(border_pt)
    else:
        sp.line.fill.background()
    return sp


def add_textbox(slide, text, left, top, width, height,
                size=11, bold=False, color=TEXT_MAIN,
                align=PP_ALIGN.LEFT, italic=False, wrap=True, font_name="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return tb


def build_title_bar(slide, title, tag_text, tag_color, slide_num, total):
    # 타이틀 배경 패널
    add_rect_shape(slide, Inches(0.3), Inches(0.2), Inches(12.733), Inches(0.85), CARD_BG, BORDER_DARK, 1.0)

    # 타이틀 텍스트
    add_textbox(slide, title, Inches(0.45), Inches(0.24), Inches(10.8), Inches(0.75), size=16, bold=True, color=TEXT_MAIN)

    # 태그 배지
    if tag_text:
        add_rect_shape(slide, Inches(11.35), Inches(0.35), Inches(1.05), Inches(0.3), RGBColor(0x0a, 0x14, 0x22), tag_color, 1.0)
        add_textbox(slide, tag_text.upper(), Inches(11.35), Inches(0.35), Inches(1.05), Inches(0.3), size=7.5, bold=True, color=tag_color, align=PP_ALIGN.CENTER)

    # 슬라이드 번호
    add_textbox(slide, f"{slide_num} / {total}", Inches(12.35), Inches(0.38), Inches(0.6), Inches(0.25), size=9, color=TEXT_SUB, align=PP_ALIGN.RIGHT)

    # 구분선
    add_rect_shape(slide, Inches(0.3), Inches(1.10), Inches(12.733), Inches(0.015), BORDER_DARK)


def build_info_card(slide, title, body, left, top, width, height, classes):
    accent = Q_COLOR
    bg = RGBColor(0x10, 0x18, 0x28)
    for cls in classes:
        if cls in CARD_COLORS:
            accent, bg = CARD_COLORS[cls]
            break

    add_rect_shape(slide, left, top, width, height, bg, accent, 1.0)
    add_rect_shape(slide, left, top, Inches(0.05), height, accent)

    cur = top + Inches(0.08)
    if title:
        add_textbox(slide, title, left + Inches(0.1), cur, width - Inches(0.18), Inches(0.32), size=10, bold=True, color=accent)
        cur += Inches(0.32)

    if body:
        body_h = height - (cur - top) - Inches(0.05)
        if body_h > Inches(0.2):
            add_textbox(slide, body, left + Inches(0.1), cur, width - Inches(0.18), body_h, size=9, color=TEXT_MAIN, wrap=True)


def build_human_note(slide, text, top):
    add_rect_shape(slide, Inches(0.3), top, Inches(12.733), Inches(0.68), RGBColor(0x1c, 0x14, 0x04), SCORE_COLOR, 1.2)
    add_rect_shape(slide, Inches(0.3), top, Inches(0.06), Inches(0.68), SCORE_COLOR)
    add_textbox(slide, f"💡 문과적 직관: {text}", Inches(0.45), top + Inches(0.06), Inches(12.4), Inches(0.56), size=9.5, color=HUMAN_NOTE_FG, wrap=True)


def build_pipeline(slide, steps, top):
    n = len(steps)
    if n == 0:
        return top
    gap = Inches(0.1)
    arr_w = Inches(0.25)
    total_arr = (n - 1) * arr_w
    total_gap = n * gap
    step_w = (Inches(12.733) - total_arr - total_gap) / n

    for i, step in enumerate(steps):
        sx = Inches(0.3) + i * (step_w + arr_w + gap)
        add_rect_shape(slide, sx, top, step_w, Inches(0.5), CARD_BG, Q_COLOR, 1.0)
        add_textbox(slide, step, sx + Inches(0.04), top + Inches(0.05), step_w - Inches(0.08), Inches(0.4), size=8.5, color=Q_COLOR, align=PP_ALIGN.CENTER, wrap=True)
        if i < n - 1:
            ax = sx + step_w + gap * 0.5
            add_textbox(slide, "→", ax, top + Inches(0.08), arr_w, Inches(0.3), size=11, color=TEXT_SUB, align=PP_ALIGN.CENTER)
    return top + Inches(0.58)


# ══════════════════════════════════════════════════════════════════════════════
#  1. 브라우저 스크린샷 캡처 (슬라이드 전체 & 수식/다이어그램 요소)
# ══════════════════════════════════════════════════════════════════════════════

def capture_all_slides_and_math(html_path: Path):
    """Playwright(Edge)를 사용하여 모든 슬라이드와 수식 요소를 고화질 캡처합니다."""
    SLIDES_IMG_DIR.mkdir(parents=True, exist_ok=True)
    MATH_IMG_DIR.mkdir(parents=True, exist_ok=True)

    print(f"[*] Playwright 실행 및 HTML 로드 중: {html_path}")
    captured_info = []

    with sync_playwright() as p:
        browser = p.chromium.launch(channel="msedge", headless=True)
        page = browser.new_page(viewport={"width": 1920, "height": 1080}, device_scale_factor=2)

        file_url = "file:///" + str(html_path.resolve()).replace("\\", "/")
        page.goto(file_url)
        page.wait_for_timeout(1500)

        # KaTeX 렌더링 강제 재실행 및 네비게이션 UI 숨기기
        page.evaluate("""() => {
            if (window.renderKaTeX) { window.renderKaTeX(); }
            const hb = document.querySelector('.header-bar');
            if (hb) hb.style.display = 'none';
            const pb = document.querySelector('.progress-bar-container');
            if (pb) pb.style.display = 'none';
        }""")
        page.wait_for_timeout(500)

        total_slides = page.locator(".slide").count()
        print(f"[*] 총 {total_slides}개 슬라이드 캡처 시작...")

        for idx in range(total_slides):
            # 슬라이드 활성화 (idx번만 표시)
            page.evaluate(f"""() => {{
                const slides = document.querySelectorAll('.slide');
                slides.forEach((s, i) => {{
                    if (i === {idx}) {{
                        s.classList.add('active');
                        s.style.display = 'flex';
                        s.style.opacity = '1';
                        s.style.transform = 'none';
                        s.style.position = 'relative';
                    }} else {{
                        s.classList.remove('active');
                        s.style.display = 'none';
                    }}
                }});
            }}""")
            page.wait_for_timeout(300)

            # 1) 슬라이드 전체 스크린샷 캡처
            slide_el = page.locator(".slide.active").first
            slide_img_path = SLIDES_IMG_DIR / f"slide_{idx+1:02d}.png"
            slide_el.screenshot(path=str(slide_img_path))

            # 2) 슬라이드 내 수식 / 다이어그램 개별 캡처 (formula-card, math-box 등)
            math_elements = []
            math_locators = page.locator(".slide.active .formula-card, .slide.active .math-box")
            math_count = math_locators.count()

            for m_idx in range(math_count):
                m_el = math_locators.nth(m_idx)
                m_path = MATH_IMG_DIR / f"slide_{idx+1:02d}_math_{m_idx+1}.png"
                m_el.screenshot(path=str(m_path))
                math_elements.append(str(m_path))

            captured_info.append({
                "slide_idx": idx,
                "slide_img": str(slide_img_path),
                "math_imgs": math_elements,
            })
            print(f"  - [{idx+1:02d}/{total_slides}] 캡처 완료 (수식 이미지 {math_count}개)")

        browser.close()

    print("[+] 모든 슬라이드 및 수식 캡처 완료!\n")
    return captured_info


# ══════════════════════════════════════════════════════════════════════════════
#  2. 고품질 PDF 생성 (100% 원본 렌더링 보존)
# ══════════════════════════════════════════════════════════════════════════════

def generate_pdf_from_captures(captured_info: list, output_path: Path):
    """캡처된 1920x1080 2x 고해상도 슬라이드 이미지들을 하나의 고품질 PDF로 병합합니다."""
    print(f"[*] PDF 생성 중: {output_path}")
    image_list = []

    for item in captured_info:
        img = Image.open(item["slide_img"]).convert("RGB")
        image_list.append(img)

    if image_list:
        first_img = image_list[0]
        first_img.save(
            str(output_path),
            save_all=True,
            append_images=image_list[1:],
            resolution=150.0,
            quality=95,
        )
        print(f"[+] PDF 생성 완료! ({len(image_list)} 페이지) -> {output_path}\n")


# ══════════════════════════════════════════════════════════════════════════════
#  3. 프레젠테이션용 풀 슬라이드 PPTX 생성 (100% 비주얼 + 발표자 노트)
# ══════════════════════════════════════════════════════════════════════════════

def generate_full_presentation_pptx(html_path: Path, captured_info: list, output_path: Path):
    """각 슬라이드의 원본 고화질 렌더링 이미지와 발표자 메모를 결합한 프레젠테이션용 PPTX를 생성합니다."""
    print(f"[*] 프레젠테이션용 풀 렌더링 PPTX 생성 중: {output_path}")
    soup = BeautifulSoup(html_path.read_text(encoding="utf-8"), "html.parser")
    slides_divs = soup.find_all("div", class_="slide")

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    for idx, (s_div, c_info) in enumerate(zip(slides_divs, captured_info)):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(slide)

        img_path = c_info["slide_img"]
        slide.shapes.add_picture(img_path, Inches(0.2), Inches(0.2), Inches(12.933), Inches(7.1))

        sn = s_div.find("div", class_="speaker-note")
        if sn:
            sn_text = clean_text(sn)
            slide.notes_slide.notes_text_frame.text = sn_text

    actual_path = safe_save_prs(prs, output_path, BASE_DIR / "transformer_seminar_presentation_new.pptx")
    print(f"[+] 풀 렌더링 PPTX 생성 완료! -> {actual_path}\n")
    return actual_path



# ══════════════════════════════════════════════════════════════════════════════
#  4. 수식 캡처 이미지 삽입형 하이브리드 PPTX 생성 (편집 가능 텍스트 + 수식 캡처)
# ══════════════════════════════════════════════════════════════════════════════

def parse_html_slides(html_path: Path):
    soup = BeautifulSoup(html_path.read_text(encoding="utf-8"), "html.parser")
    result = []

    for idx, div in enumerate(soup.find_all("div", class_="slide")):
        title_div = div.find("div", class_="slide-title")
        title = tag = ""
        tag_cls = "tag-concept"
        if title_div:
            spans = title_div.find_all("span", recursive=False)
            title = clean_text(spans[0]) if spans else ""
            if len(spans) > 1:
                tag_el = spans[1]
                tag = clean_text(tag_el)
                for c in tag_el.get("class", []):
                    if c in TAG_COLORS:
                        tag_cls = c
                        break

        slide_num = idx + 1
        body = div.find("div", class_="slide-body")
        blocks = []
        human_note = speaker_note = ""

        if body:
            intro = body.find("p", recursive=False)
            if intro:
                blocks.append({"type": "intro", "text": clean_text(intro)})

            for pf in body.find_all("div", class_="pipeline-flow"):
                steps = [clean_text(s) for s in pf.find_all("div", class_="pipeline-step")]
                if steps:
                    blocks.append({"type": "pipeline", "steps": steps})

            for card in body.find_all("div", class_="info-card"):
                if "human-note" in card.get("class", []):
                    continue
                classes = card.get("class", [])
                ct_el = card.find(class_="card-title")
                ct = clean_text(ct_el) if ct_el else ""
                card_body_el = BeautifulSoup(str(card), "html.parser").find("div")
                if card_body_el:
                    for rm in card_body_el.find_all(class_="card-title"):
                        rm.decompose()
                    cb = clean_text(card_body_el)[:400]
                else:
                    cb = ""
                blocks.append({
                    "type": "card",
                    "title": ct,
                    "body": cb,
                    "classes": classes,
                })

            hn = body.find("div", class_="human-note")
            if hn:
                human_note = clean_text(hn)
                human_note = re.sub(r"^[💡 ]*문과적 직관 해석:\s*", "", human_note)

            sn = body.find("div", class_="speaker-note")
            if sn:
                speaker_note = clean_text(sn)

        result.append({
            "idx": idx,
            "title": title,
            "tag": tag,
            "tag_cls": tag_cls,
            "slide_num": slide_num,
            "blocks": blocks,
            "human_note": human_note,
            "speaker_note": speaker_note,
        })

    return result


def generate_hybrid_pptx(html_path: Path, captured_info: list, output_path: Path):
    """수식 영역을 브라우저 캡처 이미지로 삽입하고, 나머지 텍스트/카드를 구성하는 PPTX를 생성합니다."""
    print(f"[*] 하이브리드 PPTX (수식 캡처 이미지 삽입형) 생성 중: {output_path}")
    slides_data = parse_html_slides(html_path)
    total = len(slides_data)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    for sd, c_info in zip(slides_data, captured_info):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(slide)

        tag_color = TAG_COLORS.get(sd["tag_cls"], Q_COLOR)
        build_title_bar(slide, sd["title"], sd["tag"], tag_color, sd["slide_num"], total)

        blocks = sd["blocks"]
        human = sd["human_note"]
        sp_note = sd["speaker_note"]
        math_imgs = c_info["math_imgs"]

        cur_y = Inches(1.18)
        content_floor = Inches(6.6) if human else Inches(7.0)

        # 인트로 단락
        intros = [b for b in blocks if b["type"] == "intro"]
        if intros:
            add_textbox(slide, intros[0]["text"], Inches(0.3), cur_y, Inches(12.733), Inches(0.42), size=10.5, color=TEXT_MAIN, wrap=True)
            cur_y += Inches(0.45)

        # 파이프라인
        for pb in [b for b in blocks if b["type"] == "pipeline"][:1]:
            cur_y = build_pipeline(slide, pb["steps"], cur_y)

        # 수식 캡처 이미지 삽입
        if math_imgs:
            for m_path in math_imgs:
                if os.path.exists(m_path) and cur_y + Inches(1.0) <= content_floor:
                    with Image.open(m_path) as im:
                        iw, ih = im.size
                    aspect = iw / ih if ih > 0 else 4.0
                    target_w = Inches(12.733)
                    target_h = target_w / aspect
                    if target_h > Inches(2.2):
                        target_h = Inches(2.2)
                        target_w = target_h * aspect

                    img_x = Inches(0.3) + (Inches(12.733) - target_w) / 2
                    slide.shapes.add_picture(m_path, img_x, cur_y, target_w, target_h)
                    cur_y += target_h + Inches(0.1)

        # 카드 배치 (2열 그리드)
        cards = [b for b in blocks if b["type"] == "card"]
        if cards and cur_y < content_floor:
            avail_h = content_floor - cur_y - Inches(0.08)
            n_rows = max(1, (len(cards) + 1) // 2)
            row_h = min(Inches(1.4), avail_h / n_rows - Inches(0.1))
            row_h = max(row_h, Inches(0.55))

            i = 0
            while i < len(cards) and cur_y + row_h <= content_floor:
                row = cards[i:i+2]
                n = len(row)
                cw = Inches(6.28) if n == 2 else Inches(12.733)
                for j, c in enumerate(row):
                    cx = Inches(0.3) + j * (cw + Inches(0.17))
                    build_info_card(slide, c["title"], c["body"], cx, cur_y, cw, row_h, c.get("classes", []))
                cur_y += row_h + Inches(0.1)
                i += 2

        # 문과적 직관 박스
        if human:
            note_y = max(cur_y + Inches(0.05), Inches(6.1))
            note_y = min(note_y, Inches(6.7))
            build_human_note(slide, human, note_y)

        # 발표자 메모
        if sp_note:
            slide.notes_slide.notes_text_frame.text = sp_note

    actual_path = safe_save_prs(prs, output_path, BASE_DIR / "transformer_seminar_hybrid_new.pptx")
    print(f"[+] 하이브리드 PPTX 생성 완료! -> {actual_path}\n")
    return actual_path


# ══════════════════════════════════════════════════════════════════════════════
#  메인 실행
# ══════════════════════════════════════════════════════════════════════════════

def main():
    print("=" * 70)
    print(" Transformer Seminar 자료 생성 파이프라인 시작 ")
    print("=" * 70)

    # 1. 슬라이드 및 수식 캡처 (이미 존재하면 재사용 가능)
    captured_info = capture_all_slides_and_math(HTML_FILE)

    # 2. PDF 생성
    generate_pdf_from_captures(captured_info, PDF_OUTPUT)

    # 3. 하이브리드 PPTX 생성 (수식 캡처 이미지 삽입형)
    hybrid_path = generate_hybrid_pptx(HTML_FILE, captured_info, PPTX_HYBRID_OUTPUT)

    # 4. 풀 렌더링 프레젠테이션 PPTX 생성
    full_path = generate_full_presentation_pptx(HTML_FILE, captured_info, PPTX_FULL_OUTPUT)

    print("=" * 70)
    print(" 모든 세미나 자료 생성 완료!")
    print(f" 1. PDF 문서: {PDF_OUTPUT}")
    print(f" 2. 수식 캡처 삽입 PPTX: {hybrid_path}")
    print(f" 3. 풀 렌더링 프레젠테이션 PPTX: {full_path}")
    print("=" * 70)


if __name__ == "__main__":
    main()

