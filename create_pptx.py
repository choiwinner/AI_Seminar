"""
Transformer Seminar HTML → PowerPoint(PPTX) 변환 스크립트
transformer_seminar_new.html의 슬라이드를 파싱하여 PPTX로 생성합니다.
"""

import io
import re
import sys
from pathlib import Path

from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ─── 색상 상수 (세미나 디자인 가이드 기준) ─────────────────────────────────
BG_DARK       = RGBColor(0x07, 0x0a, 0x12)   # 배경
CARD_BG       = RGBColor(0x12, 0x1a, 0x2b)   # 카드 배경
Q_COLOR       = RGBColor(0x38, 0xbd, 0xf8)   # Query - Cyan
K_COLOR       = RGBColor(0x34, 0xd3, 0x99)   # Key - Emerald
V_COLOR       = RGBColor(0xc0, 0x84, 0xfc)   # Value - Purple
SCORE_COLOR   = RGBColor(0xfb, 0xbf, 0x24)   # Score/Attention - Amber
HW_COLOR      = RGBColor(0xf4, 0x3f, 0x5e)   # H/W Memory - Rose
TEXT_MAIN     = RGBColor(0xf8, 0xfa, 0xfc)   # 메인 텍스트
TEXT_SUB      = RGBColor(0x94, 0xa3, 0xb8)   # 서브 텍스트
HUMAN_NOTE_FG = RGBColor(0xfd, 0xe6, 0x8a)   # 문과 직관 박스 텍스트
BORDER_DARK   = RGBColor(0x25, 0x35, 0x4d)   # 테두리 다크

# 카드 타입별 (악센트색, 배경색) 매핑
CARD_COLORS = {
    "card-q":     (Q_COLOR,     RGBColor(0x04, 0x13, 0x1e)),
    "card-k":     (K_COLOR,     RGBColor(0x04, 0x1c, 0x14)),
    "card-v":     (V_COLOR,     RGBColor(0x14, 0x09, 0x1e)),
    "card-score": (SCORE_COLOR, RGBColor(0x1e, 0x13, 0x04)),
    "card-hw":    (HW_COLOR,    RGBColor(0x1e, 0x04, 0x08)),
}

TAG_COLORS = {
    "tag-concept": Q_COLOR,
    "tag-math":    K_COLOR,
    "tag-hw":      HW_COLOR,
}

# ─── 슬라이드 크기 (16:9 와이드스크린) ────────────────────────────────────
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ══════════════════════════════════════════════════════════════════════════════
#  헬퍼 함수들
# ══════════════════════════════════════════════════════════════════════════════

def clean_text(el) -> str:
    """BeautifulSoup 요소에서 순수 텍스트를 추출합니다."""
    if el is None:
        return ""
    text = el.get_text(separator=" ", strip=True)
    return re.sub(r"\s+", " ", text).strip()


def add_bg(slide):
    """슬라이드 배경을 다크 컬러로 설정합니다."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK


def add_rect_shape(slide, left, top, width, height,
                   bg_color, border_color=None, border_pt=1.0):
    """사각형 도형을 추가합니다."""
    sp = slide.shapes.add_shape(1, left, top, width, height)
    sp.fill.solid()
    sp.fill.fore_color.rgb = bg_color
    if border_color:
        sp.line.color.rgb = border_color
        sp.line.width = Pt(border_pt)
    else:
        sp.line.fill.background()
    return sp


def add_textbox(slide, text, left, top, width, height,
                size=11, bold=False, color=TEXT_MAIN,
                align=PP_ALIGN.LEFT, italic=False, wrap=True):
    """텍스트 박스를 추가합니다."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return tb


def get_card_style(classes):
    """카드 클래스 목록에서 (악센트색, 배경색)을 반환합니다."""
    for cls in classes:
        if cls in CARD_COLORS:
            return CARD_COLORS[cls]
    return (Q_COLOR, RGBColor(0x10, 0x18, 0x28))


# ══════════════════════════════════════════════════════════════════════════════
#  슬라이드 컴포넌트 빌더
# ══════════════════════════════════════════════════════════════════════════════

def build_title_bar(slide, title, tag_text, tag_color, slide_num, total):
    """상단 타이틀 바를 구성합니다."""
    # 타이틀 배경 패널
    add_rect_shape(slide,
                   Inches(0.25), Inches(0.2),
                   Inches(12.83), Inches(0.88),
                   CARD_BG, BORDER_DARK, 1.0)

    # 타이틀 텍스트
    add_textbox(slide, title,
                Inches(0.4), Inches(0.25),
                Inches(11.0), Inches(0.78),
                size=17, bold=True, color=TEXT_MAIN)

    # 태그 배지
    if tag_text:
        add_rect_shape(slide,
                       Inches(11.4), Inches(0.38),
                       Inches(1.05), Inches(0.3),
                       RGBColor(0x0a, 0x14, 0x22), tag_color, 1.0)
        add_textbox(slide, tag_text.upper(),
                    Inches(11.4), Inches(0.38),
                    Inches(1.05), Inches(0.3),
                    size=7, bold=True, color=tag_color,
                    align=PP_ALIGN.CENTER)

    # 슬라이드 번호
    add_textbox(slide, f"{slide_num} / {total}",
                Inches(12.45), Inches(0.42),
                Inches(0.85), Inches(0.25),
                size=9, color=TEXT_SUB, align=PP_ALIGN.RIGHT)

    # 구분선
    add_rect_shape(slide,
                   Inches(0.25), Inches(1.12),
                   Inches(12.83), Inches(0.018),
                   BORDER_DARK)


def build_info_card(slide, title, body, left, top, width, height, classes):
    """정보 카드를 생성합니다."""
    accent, bg = get_card_style(classes)

    # 카드 배경
    add_rect_shape(slide, left, top, width, height, bg, accent, 1.2)
    # 왼쪽 강조선
    add_rect_shape(slide, left, top, Inches(0.055), height, accent)

    cur = top + Inches(0.08)
    # 제목
    if title:
        add_textbox(slide, title,
                    left + Inches(0.12), cur,
                    width - Inches(0.18), Inches(0.34),
                    size=10, bold=True, color=accent)
        cur += Inches(0.36)

    # 본문
    if body:
        body_h = height - (cur - top) - Inches(0.06)
        if body_h > Inches(0.2):
            add_textbox(slide, body,
                        left + Inches(0.12), cur,
                        width - Inches(0.18), body_h,
                        size=9, color=TEXT_MAIN, wrap=True)


def build_formula_box(slide, text, top):
    """수식/공식 박스를 생성합니다."""
    add_rect_shape(slide,
                   Inches(0.25), top,
                   Inches(12.83), Inches(0.72),
                   RGBColor(0x06, 0x10, 0x1c), Q_COLOR, 2.0)
    add_textbox(slide, text,
                Inches(0.4), top + Inches(0.08),
                Inches(12.5), Inches(0.58),
                size=11, bold=True, color=Q_COLOR,
                align=PP_ALIGN.CENTER, wrap=True)
    return top + Inches(0.78)


def build_human_note(slide, text, top):
    """💡 문과적 직관 해석 박스를 생성합니다."""
    add_rect_shape(slide,
                   Inches(0.25), top,
                   Inches(12.83), Inches(0.68),
                   RGBColor(0x1c, 0x14, 0x04), SCORE_COLOR, 1.5)
    add_rect_shape(slide, Inches(0.25), top, Inches(0.06), Inches(0.68), SCORE_COLOR)
    add_textbox(slide, f"💡 문과적 직관: {text}",
                Inches(0.4), top + Inches(0.07),
                Inches(12.5), Inches(0.56),
                size=9.5, color=HUMAN_NOTE_FG, wrap=True)


def build_pipeline(slide, steps, top):
    """파이프라인 흐름 박스들을 생성합니다."""
    n = len(steps)
    if n == 0:
        return top
    gap = Inches(0.12)
    arr_w = Inches(0.25)
    total_arr = (n - 1) * arr_w
    total_gap = n * gap
    step_w = (Inches(12.83) - total_arr - total_gap) / n

    for i, step in enumerate(steps):
        sx = Inches(0.25) + i * (step_w + arr_w + gap)
        add_rect_shape(slide, sx, top, step_w, Inches(0.55),
                       CARD_BG, Q_COLOR, 1.0)
        add_textbox(slide, step,
                    sx + Inches(0.05), top + Inches(0.06),
                    step_w - Inches(0.1), Inches(0.44),
                    size=8.5, color=Q_COLOR,
                    align=PP_ALIGN.CENTER, wrap=True)
        if i < n - 1:
            ax = sx + step_w + gap * 0.5
            add_textbox(slide, "→",
                        ax, top + Inches(0.12),
                        arr_w, Inches(0.3),
                        size=12, color=TEXT_SUB,
                        align=PP_ALIGN.CENTER)
    return top + Inches(0.62)


# ══════════════════════════════════════════════════════════════════════════════
#  HTML 파싱
# ══════════════════════════════════════════════════════════════════════════════

def parse_slides(html_path: Path):
    """HTML 파일에서 슬라이드별 데이터를 추출합니다."""
    soup = BeautifulSoup(html_path.read_text(encoding="utf-8"), "html.parser")
    result = []

    for idx, div in enumerate(soup.find_all("div", class_="slide")):
        # ── 타이틀 & 태그 추출 ──────────────────────────────────────
        title_div = div.find("div", class_="slide-title")
        title = tag = ""
        tag_cls = "tag-concept"
        if title_div:
            spans = title_div.find_all("span", recursive=False)
            title = clean_text(spans[0]) if spans else ""
            if len(spans) > 1:
                tag_el = spans[1]
                tag = clean_text(tag_el)
                for c in tag_el.get("class", []):
                    if c in TAG_COLORS:
                        tag_cls = c
                        break

        # ── 슬라이드 번호 추출 ──────────────────────────────────────
        num_el = div.find("div", class_="slide-number")
        slide_num = idx + 1  # 정수로 관리

        # ── body 내용 파싱 ──────────────────────────────────────────
        body = div.find("div", class_="slide-body")
        blocks = []
        human_note = speaker_note = ""

        if body:
            # 인트로 단락
            intro = body.find("p", recursive=False)
            if intro:
                blocks.append({"type": "intro", "text": clean_text(intro)})

            # formula-card
            for fc in body.find_all("div", class_="formula-card"):
                blocks.append({"type": "formula", "text": clean_text(fc)[:280]})

            # math-box
            for mb in body.find_all("div", class_="math-box"):
                blocks.append({"type": "math", "text": clean_text(mb)[:280]})

            # pipeline-flow
            for pf in body.find_all("div", class_="pipeline-flow"):
                steps = [clean_text(s) for s in pf.find_all("div", class_="pipeline-step")]
                if steps:
                    blocks.append({"type": "pipeline", "steps": steps})

            # info-card (human-note, speaker-note는 이미 제외)
            for card in body.find_all("div", class_="info-card"):
                if "human-note" in card.get("class", []):
                    continue
                classes = card.get("class", [])
                # card-title 복사 후 추출
                card_copy = card.__copy__() if hasattr(card, "__copy__") else BeautifulSoup(str(card), "html.parser").find("div")
                ct_el = card.find(class_="card-title")
                ct = clean_text(ct_el) if ct_el else ""
                # card-title 요소 제거 후 body 텍스트
                card_body_el = BeautifulSoup(str(card), "html.parser").find("div")
                if card_body_el:
                    for rm in card_body_el.find_all(class_="card-title"):
                        rm.decompose()
                    cb = clean_text(card_body_el)[:380]
                else:
                    cb = ""
                blocks.append({
                    "type": "card",
                    "title": ct,
                    "body": cb,
                    "classes": classes,
                })

            # human-note
            hn = body.find("div", class_="human-note")
            if hn:
                human_note = clean_text(hn)
                human_note = re.sub(r"^[💡 ]*문과적 직관 해석:\s*", "", human_note)

            # speaker-note
            sn = body.find("div", class_="speaker-note")
            if sn:
                speaker_note = clean_text(sn)

        result.append({
            "idx": idx,
            "title": title,
            "tag": tag,
            "tag_cls": tag_cls,
            "slide_num": slide_num,
            "blocks": blocks,
            "human_note": human_note,
            "speaker_note": speaker_note,
        })

    return result


# ══════════════════════════════════════════════════════════════════════════════
#  슬라이드 빌더 (메인)
# ══════════════════════════════════════════════════════════════════════════════

def build_pptx_slide(prs, data, total):
    """슬라이드 데이터 딕셔너리로 PPTX 슬라이드를 생성합니다."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 빈 레이아웃
    add_bg(slide)

    tag_color = TAG_COLORS.get(data["tag_cls"], Q_COLOR)
    build_title_bar(slide, data["title"], data["tag"],
                    tag_color, data["slide_num"], total)

    blocks    = data["blocks"]
    human     = data["human_note"]
    sp_note   = data["speaker_note"]

    # 콘텐츠 시작 Y
    cur_y = Inches(1.18)
    # 문과 노트가 있으면 하단 영역 예약
    content_floor = Inches(6.6) if human else Inches(7.0)

    # ── 인트로 단락 ────────────────────────────────────────────────
    intros = [b for b in blocks if b["type"] == "intro"]
    if intros:
        add_textbox(slide, intros[0]["text"],
                    Inches(0.25), cur_y,
                    Inches(12.83), Inches(0.48),
                    size=11, color=TEXT_MAIN, wrap=True)
        cur_y += Inches(0.5)

    # ── 수식/공식 ──────────────────────────────────────────────────
    for fb in [b for b in blocks if b["type"] == "formula"][:1]:
        cur_y = build_formula_box(slide, fb["text"], cur_y)

    for mb in [b for b in blocks if b["type"] == "math"][:1]:
        add_rect_shape(slide, Inches(0.25), cur_y, Inches(12.83), Inches(0.72),
                       RGBColor(0x06, 0x0e, 0x1a), BORDER_DARK, 1.0)
        add_textbox(slide, mb["text"],
                    Inches(0.4), cur_y + Inches(0.08),
                    Inches(12.5), Inches(0.58),
                    size=10, color=K_COLOR, align=PP_ALIGN.CENTER, wrap=True)
        cur_y += Inches(0.78)

    # ── 파이프라인 ─────────────────────────────────────────────────
    for pb in [b for b in blocks if b["type"] == "pipeline"][:1]:
        cur_y = build_pipeline(slide, pb["steps"], cur_y)

    # ── 카드 배치 (2열 그리드) ─────────────────────────────────────
    cards = [b for b in blocks if b["type"] == "card"]
    if cards:
        avail_h = content_floor - cur_y - Inches(0.08)
        n_rows = max(1, (len(cards) + 1) // 2)
        row_h = min(Inches(1.4), avail_h / n_rows - Inches(0.1))
        row_h = max(row_h, Inches(0.6))

        i = 0
        while i < len(cards) and cur_y + row_h <= content_floor:
            row = cards[i:i+2]
            n = len(row)
            cw = Inches(6.3) if n == 2 else Inches(12.83)
            for j, c in enumerate(row):
                cx = Inches(0.25) + j * (cw + Inches(0.13))
                build_info_card(slide,
                                c["title"], c["body"],
                                cx, cur_y, cw, row_h,
                                c.get("classes", []))
            cur_y += row_h + Inches(0.1)
            i += 2

    # ── 문과적 직관 해석 박스 ──────────────────────────────────────
    if human:
        note_y = max(cur_y + Inches(0.05), Inches(6.05))
        note_y = min(note_y, Inches(6.75))
        build_human_note(slide, human, note_y)

    # ── 발표자 노트 ───────────────────────────────────────────────
    if sp_note:
        slide.notes_slide.notes_text_frame.text = sp_note

    return slide


# ══════════════════════════════════════════════════════════════════════════════
#  메인 실행
# ══════════════════════════════════════════════════════════════════════════════

def main():
    # Windows 콘솔 UTF-8 출력 강제 설정
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")

    html_path = Path(r"c:\python\AI_Seminar\transformer_seminar_new.html")
    out_path  = Path(r"c:\python\AI_Seminar\transformer_seminar.pptx")

    print(f"[1/4] HTML 파싱: {html_path}")
    slides_data = parse_slides(html_path)
    total = len(slides_data)
    print(f"      슬라이드 {total}개 감지")

    print("[2/4] PPTX 초기화 (16:9 와이드스크린)")
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    print("[3/4] 슬라이드 생성 중...")
    for sd in slides_data:
        # 타이틀은 인코딩 안전하게 출력
        safe_title = sd['title'][:65].encode('utf-8', errors='replace').decode('utf-8')
        print(f"  [{sd['slide_num']:02d}/{total}] {safe_title}")
        build_pptx_slide(prs, sd, total)

    print(f"[4/4] 저장: {out_path}")
    prs.save(str(out_path))
    print(f"\n완료! -> {out_path}  ({total} slides)")


if __name__ == "__main__":
    main()
